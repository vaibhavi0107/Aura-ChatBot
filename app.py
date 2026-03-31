import streamlit as st
from google import genai
from google.genai import types
import os
import json
import requests
import smtplib
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
from PIL import Image
from io import BytesIO
from streamlit_lottie import st_lottie
from dotenv import load_dotenv

# New Advanced Libraries
from gtts import gTTS
from docx import Document
from pptx import Presentation
import urllib.parse

# Load environment variables
load_dotenv()

# --- Configuration & Setup ---
st.set_page_config(page_title="Aura: The Multi-Modal AI", page_icon="✨", layout="centered")

# --- Session State Initialization ---
if "logged_in" not in st.session_state:
    st.session_state.logged_in = False
if "username" not in st.session_state:
    st.session_state.username = ""
if "messages" not in st.session_state:
    st.session_state.messages = []
if "quick_reply" not in st.session_state:
    st.session_state.quick_reply = None
if "bot_personality" not in st.session_state:
    st.session_state.bot_personality = "Sweet 🌸"
if "enable_voice" not in st.session_state:
    st.session_state.enable_voice = False

# --- Dynamic CSS based on Personality & Login State ---
if not st.session_state.logged_in:
    # 🌟 Gen Z Cyberpunk / Neon Theme for Login
    background_color = "#0a0a0a"
    text_color = "#e0e0e0"
    accent_color = "#00ffcc"
    accent_glow = "rgba(0, 255, 204, 0.5)"
    button_border = "#ff00ff"
    button_hover_bg = "#ff00ff"
    button_hover_text = "#ffffff"
    
    custom_css = f"""
    <style>
        .stApp {{
            background: radial-gradient(circle at top right, #2b00ff, #000000 40%), radial-gradient(circle at bottom left, #ff0055, #000000 40%) !important;
            background-color: #000 !important; 
            background-attachment: fixed !important;
            color: {{text_color}} !important; 
        }}
        h1, h2, h3 {{ color: #ffffff !important; text-align: center; font-family: 'Inter', sans-serif; text-shadow: 0 0 10px #ff00ff, 0 0 20px #00ffcc; letter-spacing: 2px; }}
        div[data-testid="stForm"] {{ background: rgba(20, 20, 20, 0.6) !important; backdrop-filter: blur(16px) !important; border-radius: 20px !important; border: 1px solid rgba(255, 255, 255, 0.1) !important; padding: 30px; box-shadow: 0 8px 32px 0 rgba(0, 0, 0, 0.5); }}
        .stButton>button {{ width: 100%; border-radius: 30px; background: linear-gradient(45deg, #ff00ff, #00ffcc) !important; color: #000 !important; border: none !important; font-weight: 900; font-size: 1.1em; transition: 0.3s; box-shadow: 0 0 15px rgba(0, 255, 204, 0.5); text-transform: uppercase; }}
        .stButton>button:hover {{ transform: scale(1.03); box-shadow: 0 0 25px #ff00ff; color: #fff !important; }}
        input {{ border-radius: 10px !important; background-color: rgba(255,255,255,0.05) !important; color: #00ffcc !important; border: 1px solid rgba(255, 255, 255, 0.2) !important; }}
        input:focus {{ border-color: #ff00ff !important; box-shadow: 0 0 10px rgba(255, 0, 255, 0.5) !important; }}
        div[data-baseweb="tab"] {{ color: #00ffcc !important; font-weight: bold; font-family: monospace; font-size: 1.2em; }}
        div[aria-selected="true"] {{ color: #ff00ff !important; border-bottom-color: #ff00ff !important; text-shadow: 0 0 8px #ff00ff; }}
    </style>
    """
else:
    theme = st.session_state.bot_personality

    if theme == "Sweet 🌸":
        background_color = "#fff0f5"
        text_color = "#4d0026"
        accent_color = "#ff66b2"
        accent_glow = "rgba(255, 102, 178, 0.4)"
        bubble_bg_user = "#ff99cc"
        bubble_bg_bot = "#ffffff"
        bubble_text_bot = "#333333"
        button_border = "#ff66b2"
        button_hover_bg = "#ff66b2"
        button_hover_text = "#ffffff"
    else:
        background_color = "#1a1625"
        text_color = "#f4eaff"
        accent_color = "#bc13fe"
        accent_glow = "rgba(188, 19, 254, 0.5)"
        bubble_bg_user = "#bc13fe"
        bubble_bg_bot = "#2a223a"
        bubble_text_bot = "#f4eaff"
        button_border = "#00ffff"
        button_hover_bg = "#00ffff"
        button_hover_text = "#1a1625"

    custom_css = f"""
    <style>
        .stApp {{ background-color: {background_color} !important; color: {text_color} !important; transition: background-color 0.5s ease; }}
        h1, h2, h3 {{ color: {accent_color} !important; text-shadow: 0 0 5px {accent_color}, 0 0 10px {accent_glow}; }}
        [data-testid="stChatMessage"] {{ background-color: {bubble_bg_bot}; color: {bubble_text_bot} !important; border-radius: 15px; padding: 10px 20px; margin-bottom: 15px; border: 1px solid {accent_glow}; box-shadow: 0 0 10px {accent_glow}; }}
        [data-testid="stChatMessage"] p {{ color: {bubble_text_bot} !important; }}
        [data-testid="stChatMessage"][data-baseweb="block"]:has(svg[title="user"]) {{ background-color: {bubble_bg_user}; color: #ffffff !important; border-left: 5px solid {button_border}; }}
        [data-testid="stChatMessage"][data-baseweb="block"]:has(svg[title="user"]) p {{ color: #ffffff !important; }}
        [data-testid="stChatMessage"][data-baseweb="block"]:has(svg[title="assistant"]) {{ border-left: 5px solid {accent_color}; }}
        .stButton>button {{ width: 100%; border-radius: 20px; background-color: transparent !important; color: {button_border} !important; border: 2px solid {button_border} !important; font-weight: bold; transition: 0.3s; box-shadow: 0 0 5px {button_border}; }}
        .stButton>button:hover {{ background-color: {button_hover_bg} !important; color: {button_hover_text} !important; box-shadow: 0 0 15px {button_border}; }}
        div[data-testid="column"] .stButton>button {{ border-color: {accent_color} !important; color: {accent_color} !important; box-shadow: 0 0 5px {accent_glow}; padding: 5px 10px; font-size: 0.9em; }}
        div[data-testid="column"] .stButton>button:hover {{ background-color: {accent_color} !important; color: #fff !important; box-shadow: 0 0 15px {accent_color}; }}
        ::-webkit-scrollbar {{ width: 8px; }}
        ::-webkit-scrollbar-track {{ background: {background_color}; }}
        ::-webkit-scrollbar-thumb {{ background: {accent_color}; border-radius: 4px; }}
    </style>
    """

st.markdown(custom_css, unsafe_allow_html=True)


@st.cache_data
def load_lottieurl(url: str):
    try:
        return requests.get(url).json() if requests.get(url).status_code == 200 else None
    except:
        return None

lottie_typing = load_lottieurl("https://assets3.lottiefiles.com/packages/lf20_t2xgmroi.json")

# --- Generators for Documents & Audio ---
def generate_docx(content):
    doc = Document()
    doc.add_heading('Generated Document', 0)
    for line in content.split('\n'):
        if line.strip():
            doc.add_paragraph(line)
    bio = BytesIO()
    doc.save(bio)
    return bio.getvalue()

def generate_pptx(content):
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "AI Presentation"
    subtitle.text = "Generated by Aura"

    bullet_slide_layout = prs.slide_layouts[1]
    for slide_text in content.split('\n\n'):
        if slide_text.strip():
            slide = prs.slides.add_slide(bullet_slide_layout)
            shapes = slide.shapes
            title_shape = shapes.title
            body_shape = shapes.placeholders[1]
            lines = slide_text.split('\n')
            title_shape.text = lines[0] if len(lines) > 0 else "Slide"
            tf = body_shape.text_frame
            for line in lines[1:]:
                p = tf.add_paragraph()
                p.text = line.replace('- ', '').strip()
    
    bio = BytesIO()
    prs.save(bio)
    return bio.getvalue()

def text_to_audio(text):
    tts = gTTS(text=text, lang='en')
    bio = BytesIO()
    tts.write_to_fp(bio)
    return bio.getvalue()


# --- Database & Auth Setup ---
USERS_FILE = "users.json"
def load_users():
    if not os.path.exists(USERS_FILE):
        with open(USERS_FILE, "w") as f: json.dump({}, f)
    with open(USERS_FILE, "r") as f: return json.load(f)

def save_users(users_data):
    with open(USERS_FILE, "w") as f: json.dump(users_data, f, indent=4)

def login_page():
    # Sticker / Animation
    col1, col2, col3 = st.columns([1, 2, 1])
    with col2:
        lottie_cool = load_lottieurl("https://assets9.lottiefiles.com/packages/lf20_jcikwtux.json")
        if lottie_cool:
            st_lottie(lottie_cool, height=180, key="login_lottie")
            
    st.markdown("<h1 style='text-align: center'>✨ Welcome to Aura ✨</h1>", unsafe_allow_html=True)
    st.markdown("<p style='text-align: center; color: #a9a9a9; font-style: italic; font-size: 1.2em;'>Your Vibe-Matched AI Bestie 💅</p>", unsafe_allow_html=True)
    st.markdown("<br>", unsafe_allow_html=True)
    
    tab1, tab2 = st.tabs(["👾 Login", "🚀 Sign Up"])
    with tab1:
        with st.form("login"):
            st.markdown("### 🔑 Enter your Realm")
            username, password = st.text_input("Username"), st.text_input("Password", type="password")
            if st.form_submit_button("Let's Go 🛸"):
                users = load_users()
                if username in users and users[username]["password"] == password:
                    st.session_state.logged_in, st.session_state.username = True, username
                    st.rerun()
                else: st.error("💀 Try again bestie, wrong credentials.")
    with tab2:
        with st.form("signup"):
            st.markdown("### 🌟 Join the Vibe")
            nu, ne, np, cp = st.text_input("Username"), st.text_input("Email"), st.text_input("Password", type="password"), st.text_input("Confirm", type="password")
            if st.form_submit_button("Create Account ✨"):
                users = load_users()
                if nu not in users and np == cp and len(nu) > 2:
                    users[nu] = {"password": np, "email": ne}
                    save_users(users)
                    st.success("W! Account created 🎉")
                else: st.error("L. Check your details and try again 🤡")


# --- Main App Page ---
def main_chatbot_page():
    st.markdown("<h1 style='text-align: center'>✨ Aura AI ✨</h1>", unsafe_allow_html=True)
    
    if not st.session_state.messages:
        st.session_state.messages.append({"role": "assistant", "content": f"Hi {st.session_state.username}! I can now hear your voice, generate Word/PPT files, and create Photos! Send a prompt."})

    with st.sidebar:
        st.header("⚙️ Aura's Settings")
        new_vibe = st.radio("🎭 Personality:", ["Sweet 🌸", "Savage 😈"], index=0 if st.session_state.bot_personality == "Sweet 🌸" else 1)
        if new_vibe != st.session_state.bot_personality:
            st.session_state.bot_personality = new_vibe
            st.rerun()
            
        st.session_state.enable_voice = st.toggle("🔊 Enable Voice Output", value=st.session_state.enable_voice)
        st.markdown("---")
        
        uploaded_file = st.file_uploader("📎 Upload Document/Image Context", type=["png", "jpg", "jpeg", "pdf", "txt", "csv", "md"])
        if uploaded_file:
            st.success(f"Attached: {uploaded_file.name}")
        st.markdown("---")
        
        try:
            default_key = os.getenv("GEMINI_API_KEY", "") or st.secrets.get("GEMINI_API_KEY", "")
        except FileNotFoundError:
            default_key = os.getenv("GEMINI_API_KEY", "")
        api_key_input = st.text_input("🔑 Gemini API Key", type="password", value=default_key)
        client = genai.Client(api_key=api_key_input) if api_key_input else None
        
        if not client: st.warning("⚠️ Enter Gemini key to chat.")
        
        if st.button("🚪 Logout"):
            st.session_state.logged_in = False
            st.rerun()

    # Special Quick Replies spanning generation tools
    qr1, qr2, qr3, qr4 = st.columns(4)
    if qr1.button("📑 Gen Doc"): st.session_state.quick_reply = "Create a Word document about the history of Artificial Intelligence."
    if qr2.button("📊 Gen PPT"): st.session_state.quick_reply = "Create a PPT presentation about Machine Learning."
    if qr3.button("🎨 Gen Photo"): st.session_state.quick_reply = "Generate a photo of a futuristic cyberpunk city."
    if qr4.button("💻 Gen Code"): st.session_state.quick_reply = "Write a Python script for a simple calculator."

    # Display History
    for i, message in enumerate(st.session_state.messages):
        with st.chat_message(message["role"]):
            if "[GENERATE_DOC]" in message["content"] or "[GENERATE_PPT]" in message["content"] or "[IMAGINE]" in message["content"]:
                # The raw command is hidden from history display, we show a clean badge
                if "[GENERATE_DOC]" in message["content"]: 
                    st.info("📑 Document generated successfully!")
                    raw_text = message["content"].replace("[GENERATE_DOC]", "").strip()
                    st.download_button("Download DOCX", generate_docx(raw_text), "aura_document.docx", key=f"doc_dl_hist_{i}")
                elif "[GENERATE_PPT]" in message["content"]:
                    st.info("📊 Presentation generated successfully!")
                    raw_text = message["content"].replace("[GENERATE_PPT]", "").strip()
                    st.download_button("Download PPTX", generate_pptx(raw_text), "aura_presentation.pptx", key=f"ppt_dl_hist_{i}")
                elif "[IMAGINE]" in message["content"]:
                    image_url = message["content"].split("[IMAGINE]")[-1].strip()
                    st.image(image_url, caption="Aura's Masterpiece")
            else:
                st.markdown(message["content"])

    # Inputs
    user_text = st.chat_input("Talk to Aura...")
    audio_val = st.audio_input("🎤 Record Voice Note (Aura will instantly reply!)")
    
    prompt = None
    if st.session_state.quick_reply:
        prompt = st.session_state.quick_reply
        st.session_state.quick_reply = None
    elif user_text:
        prompt = user_text

    # Execute
    if prompt or audio_val:
        if audio_val and not prompt:
            st.chat_message("user").markdown("🎤 *You sent a Voice Message*")
            st.session_state.messages.append({"role": "user", "content": "🎤 *You sent a Voice Message*", "is_audio": True})
        elif prompt:
            st.chat_message("user").markdown(prompt)
            st.session_state.messages.append({"role": "user", "content": prompt})

        if client:
            try:
                # --- Hybrid Image Detection (Option 3 Bypassing Gemini) ---
                user_msg_lower = prompt.lower() if prompt else ""
                image_keywords = ["generate image", "create image", "draw", "photo of", "picture of", "generate a photo", "generate photo", "imagine"]
                
                is_direct_image = prompt and any(kw in user_msg_lower for kw in image_keywords)
                
                with st.chat_message("assistant"):
                    if is_direct_image:
                        st.info("🎨 Fetching Image...")
                        # Extracting keywords from prompt, ignoring common stop words
                        words = [w for w in prompt.lower().split() if w not in ["generate", "an", "image", "of", "a", "the", "create", "draw", "photo", "picture", "some"]]
                        keyword_str = ",".join(words[:4]) if words else "aesthetic"
                        
                        # Using LoremFlickr which is extremely stable and never 500s.
                        import random
                        rand_id = random.randint(1, 1000)
                        image_url = f"https://loremflickr.com/800/500/{keyword_str}?lock={rand_id}"
                        st.image(image_url, caption="Aura's Masterpiece")
                        st.session_state.messages.append({"role": "assistant", "content": f"[IMAGINE] {image_url}"})
                    else:
                        # Advanced System Prompt
                        sys_inst = f"Your name is Aura. Personality: {st.session_state.bot_personality}. Answer naturally. "
                        sys_inst += "CRITICAL INSTRUCTIONS (YOU MUST OBEY THESE): "
                        sys_inst += "1. If asked for a document/essay, you MUST start your response exactly with [GENERATE_DOC]. "
                        sys_inst += "2. If asked for a presentation/PPT, you MUST start your response exactly with [GENERATE_PPT]. Use double newlines for new slides. "
                        sys_inst += "3. If the user asks for an image, photo, drawing, or picture, YOU MUST ACT AS AN EXPERT IMAGE PROMPT GENERATOR. Output EXACTLY the tag [IMAGINE] followed by a highly detailed visual scene description. DO NOT apologize, DO NOT say you cannot generate images, and DO NOT add conversational filler (e.g., 'Here is a description:'). The system will automatically intercept your tag and render the requested image! "
                        sys_inst += "For these 3 cases, NEVER write any introductory text before the bracket tags! Just output the tag and the content. "
                        
                        hist = "\\n".join([("User: " if m["role"] == "user" else "Aura: ") + m["content"] for m in st.session_state.messages[:-1][-6:]])
                        full_prompt = f"System: {sys_inst}\\nHistory: {hist}\\nUser: {prompt if prompt else 'Attached Audio Message.'}\\nAura:"
                        
                        ph = st.empty()
                        if lottie_typing:
                            with ph.container(): st_lottie(lottie_typing, height=50)
                        else: ph.markdown("Aura is thinking...")
                            
                        contents = []
                        if uploaded_file:
                            # Pass the raw file directly to Gemini
                            file_part = types.Part.from_bytes(data=uploaded_file.getvalue(), mime_type=uploaded_file.type)
                            contents.append(file_part)
                        if audio_val:
                            audio_part = types.Part.from_bytes(data=audio_val.read(), mime_type="audio/webm")
                            contents.append(audio_part)
                        
                        contents.append(full_prompt)

                        response = client.models.generate_content(model='gemini-2.5-flash', contents=contents)
                        resp_text = response.text
                        
                        # Display real UI based on response
                        ph.empty()
                        if "[GENERATE_DOC]" in resp_text:
                            st.info("📑 Document generated!")
                            clean_text = resp_text.replace("[GENERATE_DOC]", "").strip()
                            st.download_button("Download DOCX", generate_docx(clean_text), "aura_document.docx", key="doc_dl_live")
                        elif "[GENERATE_PPT]" in resp_text:
                            st.info("📊 Presentation generated!")
                            clean_text = resp_text.replace("[GENERATE_PPT]", "").strip()
                            st.download_button("Download PPTX", generate_pptx(clean_text), "aura_presentation.pptx", key="ppt_dl_live")
                        elif "[IMAGINE]" in resp_text:
                            st.info("🎨 Generating Image...")
                            img_prompt = resp_text.split("[IMAGINE]")[-1].strip().replace('\n', ' ')
                            # Pollinations AI fails on very long prompts or newlines. Truncate strictly.
                            short_prompt = img_prompt[:150].strip()
                            safe_prompt = urllib.parse.quote(short_prompt)
                            image_url = f"https://image.pollinations.ai/prompt/{safe_prompt}?width=800&height=500&nologo=true"
                            st.image(image_url, caption="Aura's Masterpiece")
                        else:
                            st.markdown(resp_text)
                            if st.session_state.enable_voice:
                                audio_data = text_to_audio(resp_text[:300]) # Speak first 300 chars to save time
                                st.audio(audio_data, format="audio/mp3", autoplay=True)
                        
                        st.session_state.messages.append({"role": "assistant", "content": resp_text})
            except Exception as e:
                st.error(f"Error. Please ensure your Gemini Key is set. Detail: {e}")
        else: st.error("Please configure API key in sidebar.")

if st.session_state.logged_in: main_chatbot_page()
else: login_page()
